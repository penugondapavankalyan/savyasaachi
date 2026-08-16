"""Inspect reference PPTX and print full structure."""
import sys, os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'package'))

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

PPTX_PATH = os.path.join(os.path.dirname(__file__), '..', 'docs', 'reference ppt', 'savyasaachi_weekly_report.pptx')

prs = Presentation(PPTX_PATH)

print(f"Slide count : {len(prs.slides)}")
print(f"Slide size  : {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"")
print()

def emu_to_in(v): return round(v / 914400, 3)
def rgb_str(run):
    try:
        c = run.font.color
        if c and c.type:
            return str(c.rgb)
    except:
        pass
    return None

for si, slide in enumerate(prs.slides):
    print(f"{'='*60}")
    print(f"SLIDE {si+1}  layout={slide.slide_layout.name!r}")
    print(f"{'='*60}")
    for shape in slide.shapes:
        l = emu_to_in(shape.left or 0)
        t = emu_to_in(shape.top or 0)
        w = emu_to_in(shape.width or 0)
        h = emu_to_in(shape.height or 0)
        print(f"  [{shape.shape_type}] {shape.name!r}  pos=({l}\",{t}\")  size=({w}\"x{h}\")")

        if shape.has_text_frame:
            for pi, para in enumerate(shape.text_frame.paragraphs):
                txt = para.text.strip()
                if not txt:
                    continue
                sz = bold = color = align = None
                if para.runs:
                    r0 = para.runs[0]
                    sz    = round(r0.font.size.pt, 1) if r0.font.size else None
                    bold  = r0.font.bold
                    color = rgb_str(r0)
                try:
                    align = para.alignment.name if para.alignment else None
                except:
                    align = None
                print(f"    p{pi}: {txt!r}  sz={sz} bold={bold} color={color} align={align}")

        if shape.has_table:
            tbl = shape.table
            ncols = len(tbl.columns)
            nrows = len(tbl.rows)
            col_widths = [round(c.width/914400, 3) for c in tbl.columns]
            print(f"    TABLE {nrows}x{ncols}  col_widths={col_widths}")
            for ri, row in enumerate(tbl.rows):
                row_h = round(row.height/914400, 3) if row.height else None
                cells = []
                for ci, cell in enumerate(row.cells):
                    txt = cell.text.strip()
                    # get fill of cell
                    try:
                        fill = cell.fill
                        ftype = fill.type
                        if ftype and ftype.name == 'SOLID':
                            fc = str(fill.fore_color.rgb)
                        else:
                            fc = None
                    except:
                        fc = None
                    # get font of first run
                    fsz = fbold = fcolor = None
                    try:
                        tf = cell.text_frame
                        if tf.paragraphs and tf.paragraphs[0].runs:
                            r0 = tf.paragraphs[0].runs[0]
                            fsz   = round(r0.font.size.pt, 1) if r0.font.size else None
                            fbold = r0.font.bold
                            fcolor = rgb_str(r0)
                    except:
                        pass
                    cells.append(f"{txt!r}(fill={fc},fsz={fsz},bold={fbold},fc={fcolor})")
                print(f"    row{ri}(h={row_h}): {cells}")

        # Background fill of shape
        try:
            fill = shape.fill
            if fill.type and fill.type.name == 'SOLID':
                print(f"    SHAPE_FILL: #{fill.fore_color.rgb}")
        except:
            pass
    print()
