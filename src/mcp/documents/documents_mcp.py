"""
Documents MCP implementation.

No DB writes.  Generates PDF invoices (fpdf2) and PPTX analytics decks
(python-pptx) in Lambda /tmp, sends to Telegram, then deletes.
"""

from __future__ import annotations

import os
from datetime import date, timedelta
from typing import TYPE_CHECKING

from src.mcp.documents.models import AnalysisPPTXResult, InvoicePDFResult

if TYPE_CHECKING:
    from src.mcp.analytics.analytics_mcp import AnalyticsMCP
    from src.mcp.billing.billing_mcp import BillingMCP
    from src.mcp.identity.identity_mcp import IdentityMCP
    from src.mcp.payments.payments_mcp import PaymentsMCP

# ---------------------------------------------------------------------------
# State code → state name mapping (separate from system_prompt.py copy)
# ---------------------------------------------------------------------------
_STATE_NAMES: dict[str, str] = {
    "01": "Jammu & Kashmir",
    "02": "Himachal Pradesh",
    "03": "Punjab",
    "04": "Chandigarh",
    "05": "Uttarakhand",
    "06": "Haryana",
    "07": "Delhi",
    "08": "Rajasthan",
    "09": "Uttar Pradesh",
    "10": "Bihar",
    "11": "Sikkim",
    "12": "Arunachal Pradesh",
    "13": "Nagaland",
    "14": "Manipur",
    "15": "Mizoram",
    "16": "Tripura",
    "17": "Meghalaya",
    "18": "Assam",
    "19": "West Bengal",
    "20": "Jharkhand",
    "21": "Odisha",
    "22": "Chhattisgarh",
    "23": "Madhya Pradesh",
    "24": "Gujarat",
    "26": "Dadra & Nagar Haveli and Daman & Diu",
    "27": "Maharashtra",
    "28": "Andhra Pradesh",
    "29": "Karnataka",
    "30": "Goa",
    "31": "Lakshadweep",
    "32": "Kerala",
    "33": "Tamil Nadu",
    "34": "Puducherry",
    "35": "Andaman & Nicobar Islands",
    "36": "Telangana",
    "37": "Andhra Pradesh (New)",
    "38": "Ladakh",
}

# ---------------------------------------------------------------------------
# PDF Invoice Themes
#
# To add a new theme:
#   1. Add a new dict entry to _PDF_THEMES with a unique key.
#   2. Change _ACTIVE_PDF_THEME to the new key.
#
# Theme fields:
#   accent       – R,G,B  header text colour (shop name)
#   header_bg    – R,G,B  table header fill colour
#   header_fg    – R,G,B  table header text colour
#   grand_bg     – R,G,B  grand total row fill colour
#   grand_fg     – R,G,B  grand total row text colour
#   address_fg   – R,G,B  shop address / meta text colour (grey-ish)
#   divider      – R,G,B  horizontal rule colour
#   body_text    – R,G,B  regular body text colour
# ---------------------------------------------------------------------------
_PDF_THEMES: dict[str, dict] = {
    "default": {
        "accent":    (180, 140, 40),    # gold — shop name
        "header_bg": (30,  30,  30),    # near-black
        "header_fg": (255, 255, 255),   # white
        "grand_bg":  (30,  30,  30),
        "grand_fg":  (255, 255, 255),
        "address_fg": (130, 130, 130),  # grey
        "divider":   (180, 180, 180),
        "body_text": (30,  30,  30),
    },
    "blue": {
        "accent":    (13,  71, 161),    # deep blue
        "header_bg": (13,  71, 161),
        "header_fg": (255, 255, 255),
        "grand_bg":  (13,  71, 161),
        "grand_fg":  (255, 255, 255),
        "address_fg": (100, 100, 100),
        "divider":   (189, 213, 255),
        "body_text": (20,  20,  20),
    },
    "green": {
        "accent":    (27, 94, 32),
        "header_bg": (27, 94, 32),
        "header_fg": (255, 255, 255),
        "grand_bg":  (27, 94, 32),
        "grand_fg":  (255, 255, 255),
        "address_fg": (100, 130, 100),
        "divider":   (165, 214, 167),
        "body_text": (20,  20,  20),
    },
}

# ← Change this to switch themes globally
_ACTIVE_PDF_THEME = "default"

# ---------------------------------------------------------------------------
# Branding footer
# Replace the empty string with the real URL when ready, e.g.:
#   _SAVYASAACHI_URL = "https://savyasaachi.com"
# ---------------------------------------------------------------------------
# _SAVYASAACHI_URL: str = "https://savyasaachi.com"   # ← fill in when URL is ready
_SAVYASAACHI_URL: str = os.environ.get("SAVYASAACHI_INVOICE_URL", "")   # ← fill in when URL is ready

def _theme() -> dict:
    """Return the active PDF theme dict."""
    return _PDF_THEMES.get(_ACTIVE_PDF_THEME, _PDF_THEMES["default"])


# ---------------------------------------------------------------------------
# Column layout for the items table
#
# Each entry: (header_label, field_key, width_mm, align)
# field_key is used by the row renderer to pick the right value.
# ---------------------------------------------------------------------------
_ITEM_COLUMNS = [
    ("#",          "idx",            8,  "C"),
    ("ITEM",       "item",          52,  "L"),
    ("QTY",        "qty",           20,  "C"),
    ("PRICE/UNIT", "unit_price",    22,  "R"),
    ("TOTAL PRICE","taxable_value", 25,  "R"),
    ("GST %",      "gst_rate",      13,  "C"),
    ("CGST",       "cgst",          18,  "R"),
    ("SGST",       "sgst",          18,  "R"),
    ("TOTAL",      "total",         14,  "R"),
]
# Total usable width on A4 with 10mm margins each side = 190mm.
# Sum of above = 190 — verify: 8+52+20+22+25+13+18+18+14 = 190 ✓


class DocumentsMCP:
    """PDF invoice + PPTX analytics deck generation."""

    def __init__(
        self,
        billing_mcp: "BillingMCP",
        analytics_mcp: "AnalyticsMCP",
        identity_mcp: "IdentityMCP",
        payments_mcp: "PaymentsMCP | None" = None,
    ) -> None:
        self._billing = billing_mcp
        self._analytics = analytics_mcp
        self._identity = identity_mcp
        self._payments = payments_mcp

    # ------------------------------------------------------------------
    # PDF Invoice
    # ------------------------------------------------------------------

    async def generate_invoice_pdf(
        self,
        bill_id: str,
        store_id: str,
        telegram_user_id: int,
    ) -> InvoicePDFResult:
        """
        Generate a GST-correct PDF invoice for a finalized bill.
        Returns the file path; the handler sends it to Telegram.
        """
        from fpdf import FPDF  # type: ignore

        bill  = await self._billing.get_bill(bill_id)
        store = await self._identity.get_store(telegram_user_id)
        if not store:
            raise ValueError("Store not found for this user.")

        from src.utils.ist import now_ist as _now_ist
        th         = _theme()
        state_name = _STATE_NAMES.get(store.state_code, store.state_code)
        bill_date  = bill.created_at[:10] if bill.created_at else date.today().isoformat()
        gen_ts     = _now_ist().strftime("%Y-%m-%d, %H:%M:%S")

        # ── Fetch payment row (best-effort; None for PENDING_PAYMENT) ────
        pmt = None
        if self._payments:
            try:
                pmt = await self._payments.get_payment_by_bill(bill.bill_id)
            except Exception:
                pmt = None

        # ── Derive payment detail values ─────────────────────────────────
        # Determine what to show for Paid / Change / Balance lines.
        bill_status = bill.status  # CONFIRMED | CANCELLED | PENDING_PAYMENT | etc.

        if bill_status == "CANCELLED":
            paid_str    = "Rs. 0"
            change_str  = "Rs. 0"
            balance_str = "Rs. 0"
        elif bill_status == "PENDING_PAYMENT" or pmt is None:
            paid_str    = "Rs. 0"
            change_str  = "Rs. 0"
            balance_str = "Rs. 0"
        else:
            # pmt is available — derive values from payment row
            paid_str = f"Rs. {pmt.paid_amount:.2f}"

            if pmt.payment_type == "UNDERPAYMENT":
                change_str  = "Rs. 0"
                balance_str = f"Rs. {pmt.balance_due:.2f} (added to khata)"
            elif pmt.payment_type in ("OVERPAYMENT", "EXACT") and pmt.change_amount > 0:
                # Overpayment: check if change was returned (no khata_entry_id) or added to khata
                if pmt.khata_entry_id:
                    change_str = f"Rs. {pmt.change_amount:.2f} (added to khata)"
                else:
                    change_str = f"Rs. {pmt.change_amount:.2f} (change returned)"
                balance_str = "Rs. 0"
            elif pmt.payment_type == "KHATA":
                # Full credit sale — customer owes entire bill amount
                paid_str    = "Rs. 0"
                change_str  = "Rs. 0"
                balance_str = f"Rs. {bill.total_amount:.2f} (added to khata)"
            else:
                # EXACT with no change
                change_str  = "Rs. 0"
                balance_str = "Rs. 0"

        # ── Fetch customer name/phone if customer_id is present ──────────
        cust_name  = "NA"
        cust_phone = "NA"
        if bill.customer_id:
            try:
                cust_resp = (
                    self._billing.db.schema("billing")
                    .table("customers")
                    .select("name, phone")
                    .eq("id", bill.customer_id)
                    .limit(1)
                    .execute()
                )
                cust_rows = cust_resp.data or []
                if cust_rows:
                    cust_name  = cust_rows[0].get("name") or "NA"
                    cust_phone = cust_rows[0].get("phone") or "NA"
            except Exception:
                pass

        pdf = FPDF(unit="mm", format="A4")
        pdf.set_margins(left=10, top=10, right=10)
        pdf.add_page()
        pdf.set_auto_page_break(auto=True, margin=15)

        # ── Header block ────────────────────────────────────────────────
        # Left column  : shop name + address + phone + GSTIN + state
        # Right column : bill meta (Bill No / Date / Payment / Status /
        #                Paid / Change / Balance / Cust Name / Cust Phone)
        # Both columns start at the same Y, side-by-side.

        PAGE_W   = 190   # usable width (A4 210mm − 2×10mm margins)
        LEFT_W   = 110   # width of shop-info column
        RIGHT_W  = 80    # width of bill-meta column
        LABEL_W  = 32    # label cell width inside right column
        VALUE_W  = RIGHT_W - LABEL_W
        RH       = 4     # row height for right-column meta rows (mm)

        header_start_y = pdf.get_y()

        # ── Helper: emit one label+value row in the right column ─────────
        def _right_row(label: str, value: str, bold_value: bool = True) -> None:
            pdf.set_font("Helvetica", "", 8)
            pdf.set_text_color(*th["body_text"])
            pdf.cell(LABEL_W, RH, label, align="L", ln=False)
            pdf.set_font("Helvetica", "B" if bold_value else "", 8)
            pdf.cell(VALUE_W, RH, value, align="L", ln=True)

        # --- Left: shop name (large, accented) ---
        pdf.set_font("Helvetica", "B", 18)
        r, g, b = th["accent"]
        pdf.set_text_color(r, g, b)
        pdf.cell(LEFT_W, 9, store.shop_name.title(), ln=False)

        # --- Right: Bill No (taller first row to align with shop name height) ---
        pdf.set_font("Helvetica", "", 8)
        pdf.set_text_color(*th["body_text"])
        pdf.cell(LABEL_W, 9, "Bill No.", align="L", ln=False)
        pdf.set_font("Helvetica", "B", 8)
        pdf.cell(VALUE_W, 9, bill.bill_number, align="L", ln=True)

        # --- Left: address ---
        left_after_name_y = pdf.get_y()
        pdf.set_font("Helvetica", "", 8)
        pdf.set_text_color(*th["address_fg"])
        pdf.cell(LEFT_W, RH, store.address or "", ln=False)

        # --- Right: Date ---
        _right_row("Date", bill_date)

        # --- Left: phone ---
        pdf.set_font("Helvetica", "", 8)
        pdf.set_text_color(*th["address_fg"])
        pdf.cell(LEFT_W, RH, f"Phone: {store.phone}" if store.phone else "", ln=False)

        # --- Right: Payment mode ---
        _right_row("Payment", bill.payment_mode)

        # --- Left: GSTIN ---
        pdf.set_font("Helvetica", "", 8)
        pdf.set_text_color(*th["address_fg"])
        pdf.cell(LEFT_W, RH, f"GSTIN: {store.gstin}" if store.gstin else "", ln=False)

        # --- Right: Status ---
        _right_row("Status", bill_status.replace("_", " "))

        # --- Left: State ---
        pdf.set_font("Helvetica", "", 8)
        pdf.set_text_color(*th["address_fg"])
        pdf.cell(LEFT_W, RH, f"State: {state_name} ({store.state_code})", ln=False)

        # --- Right: Paid ---
        _right_row("Paid", paid_str)

        # --- Left: blank (no more left content) ---
        pdf.cell(LEFT_W, RH, "", ln=False)

        # --- Right: Change ---
        _right_row("Change", change_str)

        # --- Left: blank ---
        pdf.cell(LEFT_W, RH, "", ln=False)

        # --- Right: Balance ---
        _right_row("Balance", balance_str)

        # --- Left: blank ---
        pdf.cell(LEFT_W, RH, "", ln=False)

        # --- Right: Customer name ---
        _right_row("Customer", cust_name)

        # --- Left: blank ---
        pdf.cell(LEFT_W, RH, "", ln=False)

        # --- Right: Customer phone ---
        _right_row("Cust. Phone", cust_phone, bold_value=False)

        # --- Horizontal divider ---
        pdf.ln(2)
        pdf.set_draw_color(*th["divider"])
        pdf.set_line_width(0.5)
        pdf.line(10, pdf.get_y(), 200, pdf.get_y())
        pdf.ln(3)

        # ── Items table ──────────────────────────────────────────────────
        self._render_items_table(pdf, bill, th)

        # ── Totals block ─────────────────────────────────────────────────
        pdf.ln(4)
        self._render_totals(pdf, bill, th, PAGE_W)

        # ── Footer ───────────────────────────────────────────────────────
        pdf.ln(6)
        pdf.set_draw_color(*th["divider"])
        pdf.set_line_width(0.3)
        pdf.set_dash_pattern(dash=1, gap=2)
        pdf.line(10, pdf.get_y(), 200, pdf.get_y())
        pdf.set_dash_pattern()
        pdf.ln(3)

        pdf.set_font("Helvetica", "", 8)
        pdf.set_text_color(*th["address_fg"])
        pdf.cell(0, 5, "Thank you for your business!", align="C", ln=True)
        pdf.cell(0, 4, "This is a computer-generated invoice and does not require a signature.", align="C", ln=True)

        # ── "Powered by Savyasaachi" — inline, right-aligned, two-part ──
        # Placed immediately after the footer line above — no absolute positioning.
        # "Powered by " in very light grey; "Savyasaachi" bold + slightly darker grey.
        pdf.ln(1)
        # ── Footer last line: "Generated at …" left  |  "Powered by Savyasaachi" right ──
        # Compute widths first so both can share the same Y line.
        pdf.set_font("Helvetica", "", 7)
        _gen_label   = f"Generated at {gen_ts}"
        _gen_w       = pdf.get_string_width(_gen_label)
        _prefix      = "Powered by "
        _prefix_w    = pdf.get_string_width(_prefix)
        pdf.set_font("Helvetica", "B", 7)
        _brand_w     = pdf.get_string_width("Savyasaachi")
        _right_total = _prefix_w + _brand_w

        # "Generated at" — left margin, muted
        pdf.set_x(10)
        pdf.set_font("Helvetica", "", 7)
        pdf.set_text_color(200, 200, 200)
        pdf.cell(_gen_w, 4, _gen_label, ln=False)

        # "Powered by Savyasaachi" — flush right, same line
        pdf.set_x(200 - _right_total)
        pdf.set_text_color(200, 200, 200)
        pdf.cell(_prefix_w, 4, _prefix, ln=False)
        pdf.set_font("Helvetica", "B", 7)
        pdf.set_text_color(160, 160, 160)
        if _SAVYASAACHI_URL:
            pdf.cell(_brand_w, 4, "Savyasaachi", link=_SAVYASAACHI_URL)
        else:
            pdf.cell(_brand_w, 4, "Savyasaachi")

        import tempfile as _tempfile
        file_name = f"invoice_{bill.bill_number}.pdf"
        file_path = os.path.join(_tempfile.gettempdir(), file_name)
        pdf.output(file_path)
        file_size = os.path.getsize(file_path)

        return InvoicePDFResult(
            file_path=file_path,
            bill_number=bill.bill_number,
            file_size_bytes=file_size,
            message=f"Invoice {bill.bill_number} generated.",
        )

    # ------------------------------------------------------------------
    # PDF sub-renderers
    # ------------------------------------------------------------------

    def _render_items_table(self, pdf, bill, th: dict) -> None:
        """Render the items table with header + rows onto pdf."""
        ROW_H = 7  # mm per data row

        # ── Table header row ────────────────────────────────────────────
        r, g, b = th["header_bg"]
        pdf.set_fill_color(r, g, b)
        fr, fg, fb = th["header_fg"]
        pdf.set_text_color(fr, fg, fb)
        pdf.set_font("Helvetica", "B", 7.5)

        for label, _field, width, align in _ITEM_COLUMNS:
            pdf.cell(width, 6, label, border=0, align=align, fill=True)
        pdf.ln()

        # ── Data rows ────────────────────────────────────────────────────
        pdf.set_text_color(*th["body_text"])
        pdf.set_font("Helvetica", "", 8)

        # Alternating row fill colours
        _ROW_FILL   = (248, 245, 238)   # very light warm cream
        _ROW_PLAIN  = (255, 255, 255)   # white

        for idx, item in enumerate(bill.items, 1):
            # Determine row height — items with brand need 2 lines
            has_brand = bool(item.brand)
            row_h = ROW_H + (4 if has_brand else 0)

            # Check if row fits on current page; if not, add page + re-draw header
            if pdf.get_y() + row_h > pdf.page_break_trigger:
                pdf.add_page()
                # Re-draw table header on continuation page
                r2, g2, b2 = th["header_bg"]
                pdf.set_fill_color(r2, g2, b2)
                fr2, fg2, fb2 = th["header_fg"]
                pdf.set_text_color(fr2, fg2, fb2)
                pdf.set_font("Helvetica", "B", 7.5)
                for label, _field, width, align in _ITEM_COLUMNS:
                    pdf.cell(width, 6, label, border=0, align=align, fill=True)
                pdf.ln()
                pdf.set_text_color(*th["body_text"])
                pdf.set_font("Helvetica", "", 8)

            fill_colour = _ROW_FILL if idx % 2 == 0 else _ROW_PLAIN
            pdf.set_fill_color(*fill_colour)

            row_start_x = pdf.get_x()
            row_start_y = pdf.get_y()

            # Build values for each column
            qty_str = f"{item.quantity:.2f}\n{item.unit}"
            row_values: dict[str, str] = {
                "idx":           str(idx),
                "item":          item.product_name,
                "qty":           qty_str,
                "unit_price":    f"{item.unit_price:.2f}",
                "taxable_value": f"{item.taxable_value:.2f}",
                "gst_rate":      f"{item.gst_rate:.0f}%",
                "cgst":          f"{item.cgst_amount:.2f}",
                "sgst":          f"{item.sgst_amount:.2f}",
                "total":         f"{item.line_total:.2f}",
            }

            for col_label, field, width, align in _ITEM_COLUMNS:
                val = row_values[field]

                if field == "item":
                    # Item name: bold product name, then grey brand below
                    cell_x = pdf.get_x()
                    cell_y = pdf.get_y()
                    pdf.set_font("Helvetica", "B", 8)
                    pdf.set_fill_color(*fill_colour)
                    pdf.cell(width, ROW_H, item.product_name, border=0, align="L", fill=True)
                    if has_brand:
                        pdf.set_xy(cell_x, cell_y + ROW_H)
                        pdf.set_font("Helvetica", "", 7)
                        pdf.set_text_color(*th["address_fg"])
                        pdf.cell(width, 4, f"({item.brand})", border=0, align="L", fill=True)
                        pdf.set_text_color(*th["body_text"])
                    pdf.set_xy(cell_x + width, row_start_y)
                    pdf.set_font("Helvetica", "", 8)

                elif field == "qty":
                    # Qty: number bold on top, unit grey below
                    cell_x = pdf.get_x()
                    pdf.set_font("Helvetica", "B", 8)
                    pdf.cell(width, ROW_H, f"{item.quantity:.2f}", border=0, align="C", fill=True)
                    if has_brand:
                        pdf.set_xy(cell_x, row_start_y + ROW_H)
                        pdf.set_font("Helvetica", "", 7)
                        pdf.set_text_color(*th["address_fg"])
                        pdf.cell(width, 4, item.unit, border=0, align="C", fill=True)
                        pdf.set_text_color(*th["body_text"])
                    pdf.set_xy(cell_x + width, row_start_y)
                    pdf.set_font("Helvetica", "", 8)

                elif field == "total":
                    # Grand total column: bold
                    pdf.set_font("Helvetica", "B", 8)
                    pdf.cell(width, row_h, val, border=0, align=align, fill=True)
                    pdf.set_font("Helvetica", "", 8)

                else:
                    pdf.cell(width, row_h, val, border=0, align=align, fill=True)

            pdf.ln()

            # Thin separator line between rows (light divider colour)
            line_y = pdf.get_y()
            pdf.set_draw_color(*th["divider"])
            pdf.set_line_width(0.1)
            pdf.line(10, line_y, 200, line_y)

    def _render_totals(self, pdf, bill, th: dict, page_w: int) -> None:
        """Render the subtotal / CGST / SGST / grand-total block, right-aligned."""
        LABEL_W = 45
        VALUE_W = 35
        BLOCK_W = LABEL_W + VALUE_W
        START_X = 10 + page_w - BLOCK_W  # flush right

        pdf.set_draw_color(*th["divider"])
        pdf.set_line_width(0.2)

        rows = [
            ("Subtotal",   f"Rs. {bill.subtotal:.2f}",    False),
            ("CGST",       f"Rs. {bill.total_cgst:.2f}",  False),
            ("SGST",       f"Rs. {bill.total_sgst:.2f}",  False),
        ]
        pdf.set_font("Helvetica", "", 9)
        pdf.set_text_color(*th["body_text"])
        for label, value, _ in rows:
            pdf.set_x(START_X)
            pdf.cell(LABEL_W, 6, label, align="R", border="B")
            pdf.cell(VALUE_W, 6, value, align="R", border="B", ln=True)

        # Grand total — filled row
        r, g, b = th["grand_bg"]
        pdf.set_fill_color(r, g, b)
        fr, fg, fb = th["grand_fg"]
        pdf.set_text_color(fr, fg, fb)
        pdf.set_font("Helvetica", "B", 11)
        pdf.set_x(START_X)
        pdf.cell(LABEL_W, 9, "Grand Total", align="R", fill=True)
        pdf.cell(VALUE_W, 9, f"Rs. {bill.total_amount:.2f}", align="R", fill=True, ln=True)

    # ------------------------------------------------------------------
    # PPTX Analysis Deck
    # ------------------------------------------------------------------

    async def generate_analysis_pptx(
        self,
        store_id: str,
        telegram_user_id: int,
        period: str = "THIS_WEEK",  # TODAY | THIS_WEEK | THIS_MONTH
    ) -> AnalysisPPTXResult:
        """
        Generate a 7-slide PPTX analysis deck matching the Savyasaachi reference design.
        Returns the file path; handler sends it to Telegram and deletes it.
        """
        from pptx import Presentation                          # type: ignore
        from pptx.util import Inches, Pt, Emu                 # type: ignore
        from pptx.dml.color import RGBColor                   # type: ignore
        from pptx.enum.text import PP_ALIGN                   # type: ignore
        from pptx.chart.data import ChartData                 # type: ignore
        from pptx.enum.chart import XL_CHART_TYPE             # type: ignore

        store     = await self._identity.get_store(telegram_user_id)
        shop_name = store.shop_name.title() if store else "Your Store"

        from src.utils.ist import today_ist as _today_ist
        today = _today_ist()

        if period == "TODAY":
            start_date = end_date = today.isoformat()
            label      = f"Today ({today.strftime('%d %b %Y')})"
            date_range_label = today.strftime("%d %B %Y")
        elif period == "THIS_MONTH":
            start_date = today.replace(day=1).isoformat()
            end_date   = today.isoformat()
            label      = today.strftime("%B %Y")
            date_range_label = f"1–{today.day} {today.strftime('%B %Y')}"
        else:  # THIS_WEEK (default)
            start_date = (today - timedelta(days=today.weekday())).isoformat()
            end_date   = today.isoformat()
            label      = f"Week of {start_date}"
            from datetime import date as _d
            sd = _d.fromisoformat(start_date)
            ed = _d.fromisoformat(end_date)
            date_range_label = (
                f"{sd.day} – {ed.day} {ed.strftime('%B %Y')}"
                if sd.month == ed.month
                else f"{sd.strftime('%d %b')} – {ed.strftime('%d %b %Y')}"
            )

        data = await self._analytics.get_analytics_deck_data(
            store_id=store_id,
            store_name=shop_name,
            start_date=start_date,
            end_date=end_date,
        )

        # ── Colour palette (from reference) ──────────────────────────────
        C_TEAL      = RGBColor(0x3C, 0xB0, 0x43)   # green #3CB043 (replaces teal)
        C_ORANGE    = RGBColor(0xF4, 0xA7, 0x3D)
        C_PURPLE    = RGBColor(0x9B, 0x8F, 0xC9)
        C_DARK      = RGBColor(0x23, 0x19, 0x42)
        C_DARK2     = RGBColor(0x1F, 0x11, 0x47)
        C_MID       = RGBColor(0x74, 0x6C, 0x94)
        C_LIGHT_BG  = RGBColor(0xF5, 0xF3, 0xFB)
        C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
        C_LAVENDER  = RGBColor(0xB9, 0xAE, 0xDD)
        C_DARK_CARD = RGBColor(0x2E, 0x1D, 0x5C)
        C_RED           = RGBColor(0xD6, 0x45, 0x45)
        C_RED_BG        = RGBColor(0xFB, 0xEA, 0xEA)
        C_GREEN_STOCK   = RGBColor(0x3C, 0xB0, 0x43)   # stock-health OK dot (#3CB043)

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        blank = prs.slide_layouts[6]   # truly blank layout

        # ── Period-aggregate totals used across slides ───────────────────
        # data.summary = get_daily_summary(end_date) = today only.
        # For multi-day periods we must sum across all days in daily_summaries.
        if data.daily_summaries and len(data.daily_summaries) > 1:
            total_sales   = sum(d.total_sales   for d in data.daily_summaries)
            bill_count    = sum(d.bill_count     for d in data.daily_summaries)
            gst_collected = sum(d.total_tax      for d in data.daily_summaries)
            cash_sales    = sum(d.cash_sales     for d in data.daily_summaries)
            upi_sales     = sum(d.upi_sales      for d in data.daily_summaries)
            credit_sales  = sum(d.credit_sales   for d in data.daily_summaries)
        else:
            s = data.summary
            total_sales   = s.total_sales
            bill_count    = s.bill_count
            gst_collected = s.total_tax
            cash_sales    = s.cash_sales
            upi_sales     = s.upi_sales
            credit_sales  = s.credit_sales

        # ─────────────────────────────────────────────────────────────────
        # SLIDE 1 — Cover
        # ─────────────────────────────────────────────────────────────────
        sl1 = prs.slides.add_slide(blank)
        bg = sl1.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = C_DARK

        # decorative circles (bottom-right) — muted tones so they don't overpower the slide
        _C_CIRCLE_OUTER = RGBColor(0x2E, 0x3D, 0x6B)   # lighter navy
        _C_CIRCLE_INNER = RGBColor(0x4A, 0x4F, 0x7A)   # mid grey-purple
        _add_rect(sl1, 9.6,  3.6, 6.2, 6.2, _C_CIRCLE_OUTER, radius=True)
        _add_rect(sl1, 10.6, 4.6, 4.2, 4.2, _C_CIRCLE_INNER, radius=True)
        _add_rect(sl1, 11.5, 5.5, 2.4, 2.4, C_WHITE,          radius=True)

        # brand line — "SAVYASAACHI" stays original teal #2EC4B6 regardless of C_TEAL global
        _C_SAVYA_TEAL = RGBColor(0x2E, 0xC4, 0xB6)
        _add_text(sl1, "SAVYASAACHI", 0.6, 0.55, 6.0, 0.4,
                  size=15, bold=True, color=_C_SAVYA_TEAL)
        _add_text(sl1, "Business Intelligence Agent", 0.6, 0.95, 6.0, 0.3,
                  size=10.5, color=C_LAVENDER)
        # shop name
        _add_text(sl1, shop_name, 0.6, 2.85, 9.5, 1.1,
                  size=46, bold=True, color=C_WHITE)
        # report type
        report_type = (
            "Daily Business Report" if period == "TODAY"
            else "Monthly Business Report" if period == "THIS_MONTH"
            else "Weekly Business Report"
        )
        _add_text(sl1, report_type, 0.6, 3.78, 8.0, 0.55,
                  size=22, bold=True, color=C_ORANGE)
        # date range
        _add_text(sl1, f"Period: {date_range_label}", 1.0, 4.5, 6.0, 0.36,
                  size=14, color=C_WHITE)
        # footer
        _add_text(sl1,
                  "Generated by Savyasaachi – your always-on billing & insights agent",
                  0.6, 6.9, 9.0, 0.35, size=10.5, color=C_LAVENDER)

        # ─────────────────────────────────────────────────────────────────
        # SLIDE 2 — Week Overview  (KPI cards + payment split)
        # ─────────────────────────────────────────────────────────────────
        sl2 = prs.slides.add_slide(blank)
        _slide_header(sl2, "Week Overview",
                      f"{shop_name} – {date_range_label}",
                      C_DARK, C_MID, C_DARK2)

        # 3 KPI cards  (₹ = rupee, # = bills count, % = GST)
        _kpi_card(sl2, 0.6,  1.6, 3.843,
                  label="Total Sales",
                  value=f"Rs. {total_sales:,.2f}",
                  dot_color=C_TEAL, C_LIGHT_BG=C_LIGHT_BG,
                  C_DARK=C_DARK, C_MID=C_MID, icon="₹")
        _kpi_card(sl2, 4.743, 1.6, 3.843,
                  label="Total No. of Bills",
                  value=str(bill_count),
                  dot_color=C_ORANGE, C_LIGHT_BG=C_LIGHT_BG,
                  C_DARK=C_DARK, C_MID=C_MID, icon="#")
        _kpi_card(sl2, 8.886, 1.6, 3.843,
                  label="GST Collected",
                  value=f"Rs. {gst_collected:,.2f}",
                  dot_color=C_DARK2, C_LIGHT_BG=C_LIGHT_BG,
                  C_DARK=C_DARK, C_MID=C_MID, icon="%")

        # Payment mode split label
        _add_text(sl2, "Payment Mode Split", 0.6, 3.55, 6.0, 0.4,
                  size=17, bold=True, color=C_DARK)

        # Donut chart — Cash / UPI / Credit only
        _payment_donut(sl2, cash_sales, upi_sales, credit_sales,
                       C_TEAL, C_ORANGE, C_PURPLE)

        # Legend panel (right side)
        _payment_legend(sl2, cash_sales, upi_sales, credit_sales,
                        C_TEAL, C_ORANGE, C_PURPLE,
                        C_LIGHT_BG, C_DARK, C_MID)

        _slide_footer(sl2, C_MID)

        # ─────────────────────────────────────────────────────────────────
        # SLIDE 3 — Daily Sales Trend  (stacked bar + peak/low cards)
        # ─────────────────────────────────────────────────────────────────
        sl3 = prs.slides.add_slide(blank)
        _slide_header(sl3, "Daily Sales Trend",
                      f"Payment composition by day – {date_range_label}",
                      C_DARK, C_MID, C_DARK2)

        # Stacked bar chart: Cash / UPI / Credit per day
        _daily_stacked_bar(sl3, data.daily_summaries, C_TEAL, C_ORANGE, C_PURPLE)

        # Peak / Lowest day cards
        days_with_data = [d for d in data.daily_summaries if d.bill_count > 0]
        if days_with_data:
            peak = max(days_with_data, key=lambda d: d.total_sales)
            low  = min(days_with_data, key=lambda d: d.total_sales)
            _stat_card(sl3, 9.0, 1.65, 3.73, 2.35,
                       tag="Peak Day",
                       headline=_fmt_date_short(peak.summary_date),
                       sub=f"Rs. {peak.total_sales:,.2f} – {peak.bill_count} bills",
                       dot_color=C_TEAL, C_LIGHT_BG=C_LIGHT_BG,
                       C_DARK=C_DARK, C_MID=C_MID, icon="▲")
            _stat_card(sl3, 9.0, 4.35, 3.73, 2.35,
                       tag="Lowest Day",
                       headline=_fmt_date_short(low.summary_date),
                       sub=f"Rs. {low.total_sales:,.2f} – {low.bill_count} bills",
                       dot_color=C_RED, C_LIGHT_BG=C_LIGHT_BG,
                       C_DARK=C_DARK, C_MID=C_MID, icon="▼")

        _slide_footer(sl3, C_MID)

        # ─────────────────────────────────────────────────────────────────
        # SLIDE 4 — Top Selling Items  (bar chart + top-by-revenue/vol cards)
        # ─────────────────────────────────────────────────────────────────
        sl4 = prs.slides.add_slide(blank)
        _slide_header(sl4, "Top Selling Items",
                      f"Ranked by revenue – {label.lower()}",
                      C_DARK, C_MID, C_DARK2)

        _top_items_bar(sl4, data.top_items[:10], C_TEAL, C_ORANGE)

        if data.top_items:
            top_rev = max(data.top_items, key=lambda x: x.revenue)
            top_vol = max(data.top_items, key=lambda x: x.quantity_sold)
            _stat_card(sl4, 9.0, 1.6, 3.73, 1.75,
                       tag="Top item by revenue",
                       headline=f"{top_rev.product_name[:20]} – Rs. {top_rev.revenue:,.0f}",
                       sub="",
                       dot_color=C_TEAL, C_LIGHT_BG=C_LIGHT_BG,
                       C_DARK=C_DARK, C_MID=C_MID, icon="₹")
            _stat_card(sl4, 9.0, 3.55, 3.73, 1.75,
                       tag="Top item by volume",
                       headline=f"{top_vol.product_name[:20]} – {top_vol.quantity_sold:.0f} units",
                       sub="",
                       dot_color=C_ORANGE, C_LIGHT_BG=C_LIGHT_BG,
                       C_DARK=C_DARK, C_MID=C_MID, icon="#")

        _slide_footer(sl4, C_MID)

        # ─────────────────────────────────────────────────────────────────
        # SLIDE 5 — Stock Health  (donut + item list)
        # ─────────────────────────────────────────────────────────────────
        sh = data.stock_health
        sl5 = prs.slides.add_slide(blank)
        _slide_header(sl5,
                      "Stock Health",
                      f"Total Products: {sh.total_products}  ·  In Stock: {sh.in_stock}  ·  Low: {sh.low_stock}  ·  Out: {sh.out_of_stock}",
                      C_DARK, C_MID, C_DARK2)

        _stock_donut(sl5, sh, C_TEAL, C_ORANGE, C_RED)

        # Out-of-stock alert card
        out_items = [i for i in sh.items if i.status == "OUT_OF_STOCK"]
        if out_items:
            oi = out_items[0]
            _add_rect(sl5, 0.6, 4.75, 5.2, 1.9, C_RED_BG)
            _add_rect(sl5, 0.85, 5.0, 0.55, 0.55, C_RED, radius=True)
            # Warning icon centred on the red dot
            _add_icon_on_dot(sl5, "!", dot_l=0.85, dot_t=5.0,
                             dot_size=0.55, font_size=16)
            _add_text(sl5, "Out of Stock", 1.6, 4.95, 4.0, 0.3,
                      size=11.5, bold=True, color=C_RED)
            _add_text(sl5, oi.product_name, 1.6, 5.22, 4.0, 0.4,
                      size=18, bold=True, color=C_DARK)
            _add_text(sl5,
                      f"Current: {oi.quantity_in_stock} {oi.unit}  ·  Reorder level: {oi.reorder_level} {oi.unit}",
                      1.6, 5.63, 4.0, 0.6, size=11, color=C_DARK)

        # Stock item list (right two columns) — includes LOW/CRITICAL items with coloured dots
        ok_items = [i for i in sh.items if i.status != "OUT_OF_STOCK"]
        _stock_item_list(sl5, ok_items[:12], C_GREEN_STOCK, C_ORANGE, C_RED, C_DARK, C_MID)

        _slide_footer(sl5, C_MID)

        # ─────────────────────────────────────────────────────────────────
        # SLIDE 6 — Khata (Credit) Overview
        # ─────────────────────────────────────────────────────────────────
        ko = data.khata_overview
        sl6 = prs.slides.add_slide(blank)
        _slide_header(sl6, "Khata Overview",
                      f"Credit ledger summary – {date_range_label}",
                      C_DARK, C_MID, C_DARK2)

        # ── Two KPI cards (top-left): total credit given + total shop owes ──
        # Card 1: Total Credit Given (shop is owed) — orange dot, ₹ icon
        _kpi_card(sl6, 0.6, 1.6, 4.0,
                  label=f"Credit Given  ({ko.credit_customer_count} customers)",
                  value=f"Rs. {ko.total_credit_given:,.2f}",
                  dot_color=C_ORANGE, C_LIGHT_BG=C_LIGHT_BG,
                  C_DARK=C_DARK, C_MID=C_MID, icon="₹")

        # Card 2: Shop Owes Customers — teal/green dot, ₹ icon
        C_GREEN     = RGBColor(0x27, 0xAE, 0x60)
        C_GREEN_BG  = RGBColor(0xE8, 0xF8, 0xEE)
        _kpi_card(sl6, 4.77, 1.6, 4.0,
                  label=f"Shop Owes  ({ko.shop_owes_customer_count} customers)",
                  value=f"Rs. {ko.total_shop_owes:,.2f}",
                  dot_color=C_GREEN, C_LIGHT_BG=C_LIGHT_BG,
                  C_DARK=C_DARK, C_MID=C_MID, icon="₹")

        # ── Highest debtor card (ref image style: red bg, ! dot, bold name) ──
        if ko.highest_debtor:
            hd = ko.highest_debtor
            _add_rect(sl6, 0.6, 3.55, 4.0, 1.9, C_RED_BG)
            _add_rect(sl6, 0.85, 3.8, 0.55, 0.55, C_RED, radius=True)
            _add_icon_on_dot(sl6, "₹", dot_l=0.85, dot_t=3.8, dot_size=0.55, font_size=14)
            _add_text(sl6, "Highest Credit", 1.6, 3.75, 3.0, 0.3,
                      size=11.5, bold=True, color=C_RED)
            _add_text(sl6, hd.name[:22], 1.6, 4.03, 3.0, 0.4,
                      size=18, bold=True, color=C_DARK)
            _add_text(sl6,
                      f"Owes Rs. {hd.balance:,.2f}  ·  {hd.phone}",
                      1.6, 4.45, 3.0, 0.4, size=11, color=C_DARK)

        # ── Highest creditor card (light-green style, shop owes this customer) ──
        if ko.highest_creditor:
            hc = ko.highest_creditor
            _add_rect(sl6, 4.77, 3.55, 4.0, 1.9, C_GREEN_BG)
            _add_rect(sl6, 5.02, 3.8, 0.55, 0.55, C_GREEN, radius=True)
            _add_icon_on_dot(sl6, "₹", dot_l=5.02, dot_t=3.8, dot_size=0.55, font_size=14)
            _add_text(sl6, "We Owe Most", 5.77, 3.75, 3.0, 0.3,
                      size=11.5, bold=True, color=C_GREEN)
            _add_text(sl6, hc.name[:22], 5.77, 4.03, 3.0, 0.4,
                      size=18, bold=True, color=C_DARK)
            _add_text(sl6,
                      f"We owe Rs. {abs(hc.balance):,.2f}  ·  {hc.phone}",
                      5.77, 4.45, 3.0, 0.4, size=11, color=C_DARK)

        # ── Credit by Day bar chart (Slide-3 style, right panel) ────────
        _add_text(sl6, "Credit Sales by Day", 9.05, 1.45, 4.0, 0.4,
                  size=15, bold=True, color=C_DARK)
        _credit_by_day_bar(sl6, ko.credit_by_day, C_PURPLE)

        _slide_footer(sl6, C_MID)

        # ─────────────────────────────────────────────────────────────────
        # SLIDE 7 — GST Summary  (dark bg, KPI cards, bar chart + slab table)
        # ─────────────────────────────────────────────────────────────────
        gs = data.gst_summary
        sl6 = prs.slides.add_slide(blank)
        bg6 = sl6.background
        bg6.fill.solid()
        bg6.fill.fore_color.rgb = C_DARK

        _slide_header(sl6, "GST Summary",
                      f"Period: {date_range_label}",
                      C_WHITE, C_LAVENDER, C_TEAL,
                      savya_color=C_TEAL)

        # 4 KPI cards (dark style) — icons: ₹ / C / S / %
        gst_cards = [
            ("Total Taxable", f"Rs. {gs.total_taxable_value:,.2f}", C_TEAL,   "₹"),
            ("Total CGST",    f"Rs. {gs.total_cgst:,.2f}",          C_TEAL,   "C%"),
            ("Total SGST",    f"Rs. {gs.total_sgst:,.2f}",          C_TEAL,   "S%"),
            ("Total GST",     f"Rs. {gs.total_gst:,.2f}",           C_ORANGE, "%"),
        ]
        card_w = 2.845
        for ci, (lbl, val, dot, ico) in enumerate(gst_cards):
            cx = 0.6 + ci * (card_w + 0.25)
            _kpi_card_dark(sl6, cx, 1.6, card_w,
                           label=lbl, value=val, dot_color=dot,
                           C_DARK_CARD=C_DARK_CARD,
                           C_WHITE=C_WHITE, C_LAVENDER=C_LAVENDER,
                           icon=ico)

        # White panel: bar chart
        _add_rect(sl6, 0.6, 3.45, 7.1, 3.55, C_WHITE)
        _add_text(sl6, "GST Collected by Slab", 0.9, 3.65, 6.5, 0.35,
                  size=15, bold=True, color=C_DARK)
        _gst_bar(sl6, gs, C_TEAL, C_ORANGE)

        # Dark panel: slab breakdown text
        _add_rect(sl6, 7.95, 3.45, 4.78, 3.55, C_DARK_CARD)
        _add_text(sl6, "Slab Breakdown", 8.25, 3.65, 4.2, 0.35,
                  size=15, bold=True, color=C_WHITE)
        ty = 4.15
        for slab in gs.by_slab:
            _add_text(sl6, f"{slab.gst_rate:.0f}% Slab",
                      8.25, ty, 1.3, 0.3, size=12, bold=True, color=C_ORANGE)
            _add_text(sl6,
                      f"Taxable Rs. {slab.taxable_value:,.2f}  ·  GST Rs. {slab.total_gst:,.2f}  ·  {slab.item_count} items",
                      8.25, ty + 0.28, 4.3, 0.3, size=10.5, color=C_LAVENDER)
            # divider line
            _add_rect(sl6, 8.25, ty + 0.64, 4.2, 0.01, C_MID)
            ty += 0.78

        # footer
        _add_text(sl6, "Thank you – Powered by Savyasaachi",
                  0.6, 7.05, 6.0, 0.32, size=10.5, color=C_LAVENDER)

        import tempfile as _tempfile
        file_name = f"analysis_{store_id[:8]}_{end_date}.pptx"
        file_path = os.path.join(_tempfile.gettempdir(), file_name)
        prs.save(file_path)
        file_size = os.path.getsize(file_path)

        return AnalysisPPTXResult(
            file_path=file_path,
            period_label=label,
            file_size_bytes=file_size,
            message=f"Analysis deck for {label} generated.",
        )


# ---------------------------------------------------------------------------
# PPTX shape / text helpers
# ---------------------------------------------------------------------------

def _add_text(slide, text: str, l: float, t: float, w: float, h: float,
              size: float = 11, bold: bool = False, color=None,
              align=None) -> None:
    """Add a plain textbox."""
    from pptx.util import Inches, Pt  # type: ignore
    from pptx.enum.text import PP_ALIGN  # type: ignore
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = False
    para = tf.paragraphs[0]
    run  = para.add_run()
    run.text       = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    if color:
        run.font.color.rgb = color
    if align:
        para.alignment = align


def _add_rect(slide, l: float, t: float, w: float, h: float, fill_color,
              radius: bool = False) -> None:
    """Add a filled rectangle or oval (radius=True) with no border."""
    from pptx.util import Inches  # type: ignore
    # MSO_AUTO_SHAPE_TYPE: 1=RECTANGLE, 9=OVAL
    shape_type = 9 if radius else 1
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()   # no border


def _slide_header(slide, title: str, subtitle: str,
                  title_color, sub_color, savya_bg_color,
                  savya_color=None) -> None:
    """Standard slide header: large title, subtitle, SAVYASAACHI brand top-right."""
    from pptx.util import Inches, Pt  # type: ignore
    from pptx.enum.text import PP_ALIGN  # type: ignore
    from pptx.dml.color import RGBColor  # type: ignore
    _add_text(slide, title,    0.6, 0.38, 8.5, 0.6,  size=32, bold=True, color=title_color)
    _add_text(slide, subtitle, 0.6, 0.95, 8.5, 0.35, size=12.5, color=sub_color)
    # "SAVYASAACHI" top-right
    tb = slide.shapes.add_textbox(Inches(9.53), Inches(0.42), Inches(3.2), Inches(0.3))
    tf = tb.text_frame
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.RIGHT
    run = para.add_run()
    run.text = "SAVYASAACHI"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = savya_color if savya_color else savya_bg_color


def _slide_footer(slide, color) -> None:
    """Bottom-right 'Savyasaachi – Business Agent' footer."""
    from pptx.util import Inches, Pt  # type: ignore
    from pptx.enum.text import PP_ALIGN  # type: ignore
    tb = slide.shapes.add_textbox(Inches(9.53), Inches(7.08), Inches(3.2), Inches(0.3))
    tf = tb.text_frame
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.RIGHT
    run = para.add_run()
    run.text = "Savyasaachi – Business Agent"
    run.font.size = Pt(8.5)
    run.font.color.rgb = color


def _add_icon_on_dot(slide, icon: str, dot_l: float, dot_t: float,
                     dot_size: float = 0.5, font_size: float = 14) -> None:
    """Render an icon symbol perfectly centred (H + V) over a circle dot.

    Uses PP_ALIGN.CENTER for horizontal and MSO_ANCHOR.MIDDLE for vertical,
    with the textbox exactly matching the dot's bounding box so the symbol
    sits dead-centre regardless of character width.
    """
    from pptx.util import Inches, Pt                      # type: ignore
    from pptx.dml.color import RGBColor                   # type: ignore
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR       # type: ignore

    tb = slide.shapes.add_textbox(
        Inches(dot_l), Inches(dot_t), Inches(dot_size), Inches(dot_size)
    )
    tf = tb.text_frame
    tf.word_wrap         = False
    tf.auto_size         = None          # no auto-resize — fixed box = dot size
    tf.vertical_anchor   = MSO_ANCHOR.MIDDLE   # vertical centre

    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER     # horizontal centre

    run = para.add_run()
    run.text           = icon
    run.font.size      = Pt(font_size)
    run.font.bold      = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def _kpi_card(slide, l: float, t: float, w: float,
              label: str, value: str, dot_color,
              C_LIGHT_BG, C_DARK, C_MID,
              icon: str = "●") -> None:
    """Light KPI card: bg rect + coloured circle + centred icon + value + label."""
    _add_rect(slide, l, t, w, 1.65, C_LIGHT_BG)
    # Coloured circle
    _add_rect(slide, l + 0.22, t + 0.2, 0.5, 0.5, dot_color, radius=True)
    # Icon perfectly centred on the circle
    _add_icon_on_dot(slide, icon, dot_l=l + 0.22, dot_t=t + 0.2)
    _add_text(slide, value, l + 0.22, t + 0.78, w - 0.42, 0.42,
              size=22, bold=True, color=C_DARK)
    _add_text(slide, label, l + 0.22, t + 1.31, w - 0.42, 0.28,
              size=11, color=C_MID)


def _kpi_card_dark(slide, l: float, t: float, w: float,
                   label: str, value: str, dot_color,
                   C_DARK_CARD, C_WHITE, C_LAVENDER,
                   icon: str = "●") -> None:
    """Dark KPI card for GST slide."""
    _add_rect(slide, l, t, w, 1.6, C_DARK_CARD)
    _add_rect(slide, l + 0.22, t + 0.2, 0.5, 0.5, dot_color, radius=True)
    _add_icon_on_dot(slide, icon, dot_l=l + 0.22, dot_t=t + 0.2)
    _add_text(slide, value, l + 0.22, t + 0.78, w - 0.42, 0.42,
              size=18, bold=True, color=C_WHITE)
    _add_text(slide, label, l + 0.22, t + 1.26, w - 0.42, 0.28,
              size=11, color=C_LAVENDER)


def _stat_card(slide, l: float, t: float, w: float, h: float,
               tag: str, headline: str, sub: str, dot_color,
               C_LIGHT_BG, C_DARK, C_MID,
               icon: str = "●") -> None:
    """Info card: bg + dot + centred icon + tag label + headline + sub."""
    _add_rect(slide, l, t, w, h, C_LIGHT_BG)
    _add_rect(slide, l + 0.25, t + 0.25, 0.55, 0.55, dot_color, radius=True)
    # Icon perfectly centred on the circle
    _add_icon_on_dot(slide, icon, dot_l=l + 0.25, dot_t=t + 0.25,
                     dot_size=0.55, font_size=14)
    _add_text(slide, tag,      l + 0.25, t + 0.95, w - 0.5, 0.3,  size=12, color=C_MID)
    _add_text(slide, headline, l + 0.25, t + 1.23, w - 0.5, 0.42, size=20, bold=True, color=C_DARK)
    if sub:
        _add_text(slide, sub,  l + 0.25, t + 1.68, w - 0.5, 0.35, size=12, color=C_DARK)


def _payment_donut(slide, cash: float, upi: float, credit: float,
                   C_TEAL, C_ORANGE, C_PURPLE) -> None:
    """Donut chart: Cash / UPI / Credit.
    When all values are 0 uses equal placeholders so chart still renders.
    """
    from pptx.chart.data import ChartData   # type: ignore
    from pptx.enum.chart import XL_CHART_TYPE  # type: ignore
    from pptx.util import Inches, Pt           # type: ignore
    from pptx.dml.color import RGBColor        # type: ignore

    # Guard: all-zero → equal placeholders so chart renders as a grey ring
    if cash == 0 and upi == 0 and credit == 0:
        cash = upi = credit = 1.0

    cd = ChartData()
    cd.categories = ["Cash", "UPI", "Credit"]
    cd.add_series("Amount", (cash, upi, credit))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        Inches(0.3), Inches(3.95), Inches(5.6), Inches(3.2),
        cd,
    ).chart
    chart.has_legend = False
    chart.has_title  = False
    series = chart.series[0]
    colors = [C_TEAL, C_ORANGE, C_PURPLE]
    for i, pt in enumerate(series.points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = colors[i % len(colors)]

    # Data labels: show category name + percentage on each slice
    series.has_data_labels = True
    dl = series.data_labels
    dl.show_category_name = True
    dl.show_percentage    = True
    dl.show_value         = False
    dl.show_series_name   = False
    dl.show_legend_key    = False
    dl.font.size          = Pt(11)
    dl.font.bold          = True
    dl.font.color.rgb     = RGBColor(0xFF, 0xFF, 0xFF)


def _payment_legend(slide, cash: float, upi: float, credit: float,
                    C_TEAL, C_ORANGE, C_PURPLE,
                    C_LIGHT_BG, C_DARK, C_MID) -> None:
    """Right-side legend panel for payment split — height sized exactly to content."""
    items = [
        (C_TEAL,   "Cash",   cash),
        (C_ORANGE, "UPI",    upi),
        (C_PURPLE, "Credit", credit),
    ]
    panel_top    = 3.93
    row_h        = 0.72
    top_pad      = 0.15
    bottom_pad   = 0.15
    panel_height = top_pad + len(items) * row_h + bottom_pad
    _add_rect(slide, 6.9, panel_top, 5.83, panel_height, C_LIGHT_BG)
    y = panel_top + top_pad
    for idx, (dot_c, lbl, val) in enumerate(items):
        _add_rect(slide, 7.2, y + 0.14, 0.16, 0.16, dot_c, radius=True)
        _add_text(slide, lbl,               7.5,   y,      2.5, 0.5, size=13.5, color=C_DARK)
        _add_text(slide, f"Rs. {val:,.2f}", 10.03, y,      2.4, 0.5, size=13.5, bold=True,
                  color=C_DARK, align=None)
        # divider only between rows, not after the last one
        if idx < len(items) - 1:
            _add_rect(slide, 7.2, y + 0.62, 5.23, 0.01, C_MID)
        y += row_h


def _fmt_axis_date(date_str: str) -> str:
    """'2026-08-15' → '15th Aug'  (ordinal suffix, 3-letter month, cross-platform)"""
    from datetime import date as _d
    import sys as _sys
    try:
        d = _d.fromisoformat(date_str)
        day = d.day
        # ordinal suffix
        if 11 <= day <= 13:
            suffix = "th"
        else:
            suffix = {1: "st", 2: "nd", 3: "rd"}.get(day % 10, "th")
        return f"{day}{suffix} {d.strftime('%b')}"
    except Exception:
        return date_str


def _daily_stacked_bar(slide, daily_summaries, C_TEAL, C_ORANGE, C_PURPLE) -> None:
    """Stacked bar chart: Cash / UPI / Credit per calendar day."""
    from pptx.chart.data import ChartData          # type: ignore
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION  # type: ignore
    from pptx.util import Inches, Pt               # type: ignore
    from pptx.dml.color import RGBColor            # type: ignore

    cd = ChartData()
    cd.categories = [_fmt_axis_date(d.summary_date) for d in daily_summaries]
    cd.add_series("Cash",   tuple(d.cash_sales   for d in daily_summaries))
    cd.add_series("UPI",    tuple(d.upi_sales    for d in daily_summaries))
    cd.add_series("Credit", tuple(d.credit_sales for d in daily_summaries))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_STACKED,
        Inches(0.4), Inches(1.65), Inches(8.3), Inches(5.1),
        cd,
    ).chart
    chart.has_legend          = True
    chart.has_title           = False
    chart.legend.position     = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    colors = [C_TEAL, C_ORANGE, C_PURPLE]
    for si, series in enumerate(chart.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = colors[si]


def _credit_by_day_bar(slide, credit_by_day: list, C_PURPLE) -> None:
    """Single-series vertical bar chart: credit sales amount per day.
    Positioned on the right side of the Khata Overview slide.
    Uses _fmt_axis_date for axis labels (e.g. '15th Aug').
    """
    from pptx.chart.data import ChartData          # type: ignore
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION  # type: ignore
    from pptx.util import Inches                   # type: ignore
    from pptx.dml.color import RGBColor            # type: ignore

    cd = ChartData()
    cd.categories = [_fmt_axis_date(d) for d, _ in credit_by_day]
    cd.add_series("Credit", tuple(v for _, v in credit_by_day))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(8.9), Inches(1.85), Inches(4.23), Inches(4.55),
        cd,
    ).chart
    chart.has_legend = False
    chart.has_title  = False
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = C_PURPLE


def _top_items_bar(slide, items, C_TEAL, C_ORANGE) -> None:
    """Horizontal bar chart: top items by revenue."""
    from pptx.chart.data import ChartData   # type: ignore
    from pptx.enum.chart import XL_CHART_TYPE  # type: ignore
    from pptx.util import Inches              # type: ignore

    if not items:
        return
    cd = ChartData()
    cd.categories = [i.product_name[:20] for i in items]
    cd.add_series("Revenue", tuple(i.revenue for i in items))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.3), Inches(1.6), Inches(8.3), Inches(5.3),
        cd,
    ).chart
    chart.has_legend = False
    chart.has_title  = False
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = C_TEAL


def _stock_donut(slide, sh, C_TEAL, C_ORANGE, C_RED) -> None:
    """Donut chart: In Stock / Low / Out of Stock — with data labels showing value + %."""
    from pptx.chart.data import ChartData   # type: ignore
    from pptx.enum.chart import XL_CHART_TYPE  # type: ignore
    from pptx.util import Inches              # type: ignore
    from pptx.dml.color import RGBColor       # type: ignore

    # Guard: no stock data at all → render a placeholder
    in_s = sh.in_stock or 0
    low  = sh.low_stock or 0
    out  = sh.out_of_stock or 0
    if in_s == 0 and low == 0 and out == 0:
        in_s = 1

    cd = ChartData()
    # Only include categories that have non-zero counts so legend is clean
    categories, counts, colors_list = [], [], []
    for label, count, color in [
        ("In Stock",     in_s, C_TEAL),
        ("Low",          low,  C_ORANGE),
        ("Out of Stock", out,  C_RED),
    ]:
        if count > 0:
            categories.append(label)
            counts.append(count)
            colors_list.append(color)
    cd.categories = categories
    cd.add_series("Count", tuple(counts))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        Inches(0.3), Inches(1.55), Inches(5.6), Inches(3.0),
        cd,
    ).chart
    chart.has_legend = True
    chart.has_title  = False

    # Data labels: show percentage + value
    plot = chart.plots[0]
    plot.has_data_labels = True
    dls = plot.data_labels
    try:
        dls.show_percentage  = True
        dls.show_value       = True
        dls.show_category_name = False
        dls.show_series_name   = False
    except Exception:
        pass  # older python-pptx versions may not support all flags

    for i, pt in enumerate(chart.series[0].points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = colors_list[i % len(colors_list)]


def _stock_item_list(slide, items, C_GREEN_STOCK, C_ORANGE, C_RED, C_DARK, C_MID) -> None:
    """Two-column item list for stock slide.

    Dot colour per item status:
      OK              → C_GREEN_STOCK (#3CB043)
      LOW / CRITICAL  → C_ORANGE (amber)
      OUT_OF_STOCK    → C_RED (should not appear here, but guarded anyway)

    Item label format: "name (brand)" when brand is present, else "name".
    """
    if not items:
        return

    def _dot_color(status: str):
        if status in ("LOW", "CRITICAL"):
            return C_ORANGE
        if status == "OUT_OF_STOCK":
            return C_RED
        return C_GREEN_STOCK

    def _item_label(item) -> str:
        name = item.product_name or ""
        brand = (item.brand or "").strip()
        if brand:
            label = f"{name} ({brand.title()})"
        else:
            label = name
        return label[:26]

    mid = (len(items) + 1) // 2
    col_xs = [6.1, 9.15]
    for col_i, col_x in enumerate(col_xs):
        col_items = items[col_i * mid: (col_i + 1) * mid]
        y = 1.67
        for item in col_items:
            dot_color = _dot_color(item.status)
            _add_rect(slide, col_x, y + 0.04, 0.13, 0.13, dot_color, radius=True)
            _add_text(slide, _item_label(item), col_x + 0.24, y - 0.15,
                      2.63, 0.3, size=11.5, bold=True, color=C_DARK)
            _add_text(slide,
                      f"Stock {item.quantity_in_stock} {item.unit}  ·  Reorder {item.reorder_level}",
                      col_x + 0.24, y + 0.14, 2.63, 0.3, size=9, color=C_MID)
            # divider
            _add_rect(slide, col_x, y + 0.54, 2.87, 0.01, C_MID)
            y += 0.76


def _gst_bar(slide, gs, C_TEAL, C_ORANGE) -> None:
    """Horizontal bar chart: GST by slab."""
    from pptx.chart.data import ChartData   # type: ignore
    from pptx.enum.chart import XL_CHART_TYPE  # type: ignore
    from pptx.util import Inches              # type: ignore

    if not gs.by_slab:
        return
    cd = ChartData()
    cd.categories = [f"{int(s.gst_rate)}%" for s in gs.by_slab]
    cd.add_series("GST", tuple(s.total_gst for s in gs.by_slab))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.85), Inches(4.05), Inches(6.6), Inches(2.8),
        cd,
    ).chart
    chart.has_legend = False
    chart.has_title  = False
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = C_TEAL


def _fmt_date_short(date_str: str) -> str:
    """'2026-08-15' → '15 August'  (cross-platform: no leading zero)"""
    from datetime import date as _d
    import sys as _sys
    try:
        d = _d.fromisoformat(date_str)
        # %-d is Linux-only; %#d is Windows; fall back to lstrip("0")
        try:
            day = d.strftime("%#d") if _sys.platform == "win32" else d.strftime("%-d")
        except ValueError:
            day = str(d.day)
        return f"{day} {d.strftime('%B')}"
    except Exception:
        return date_str
